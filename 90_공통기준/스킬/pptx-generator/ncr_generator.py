"""
NCR (Non-Conformity Report) PPT 생성기 — v2
GPT 합의 기반: python-pptx + Beautify 규칙 엔진

입력: NCR 데이터 dict (photo_path 옵션)
출력: 2슬라이드 NCR 보고서 PPTX

사용법:
  PYTHONUTF8=1 python ncr_generator.py
  PYTHONUTF8=1 python ncr_generator.py --input data.json --output report.pptx
"""

import json
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Beautify 색상 팔레트 (제조업 보고서) ──
COLORS = {
    "primary":    RGBColor(0x1E, 0x27, 0x61),  # 진한 남색
    "bg_light":   RGBColor(0xF2, 0xF2, 0xF2),  # 연한 회색
    "accent":     RGBColor(0xE8, 0x70, 0x2A),  # 주황
    "text_dark":  RGBColor(0x33, 0x33, 0x33),  # 진회색
    "text_white": RGBColor(0xFF, 0xFF, 0xFF),  # 흰색
    "card_bg":    RGBColor(0xFF, 0xFF, 0xFF),  # 카드 배경
    "card_label": RGBColor(0xE8, 0xEB, 0xF0),  # 라벨 배경
    "pass_green": RGBColor(0x2C, 0x5F, 0x2D),  # 양호
    "fail_red":   RGBColor(0xC0, 0x39, 0x2B),  # 불량
}

# ── 슬라이드 크기 (16:9) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Beautify 여백/크기 규칙 ──
MARGIN = Inches(0.5)
HEADER_H = Inches(0.8)
CARD_GAP = Inches(0.08)
CARD_RADIUS = Inches(0.08)
FONT_TITLE = Pt(28)
FONT_SUBTITLE = Pt(14)
FONT_LABEL = Pt(10)
FONT_VALUE = Pt(11)
FONT_BODY = Pt(11)
FONT_SECTION = Pt(12)
FONT_PAGE = Pt(10)


def _set_run(paragraph, text, font_size, color, bold=False, align=None):
    """paragraph에 run 단위로 텍스트 설정 (QA에서 run.font.size 감지 가능)"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    if align is not None:
        paragraph.alignment = align


def add_header(slide, title, subtitle, page_text):
    """슬라이드 상단 헤더 바 추가"""
    # 헤더 배경
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W, HEADER_H  # MSO_SHAPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()

    # 제목 텍스트
    txbox = slide.shapes.add_textbox(MARGIN, Inches(0.08), Inches(8), Inches(0.4))
    tf = txbox.text_frame
    tf.word_wrap = True
    _set_run(tf.paragraphs[0], title, FONT_TITLE, COLORS["text_white"], bold=True)

    # 부제목
    txbox2 = slide.shapes.add_textbox(MARGIN, Inches(0.48), Inches(10), Inches(0.25))
    tf2 = txbox2.text_frame
    _set_run(tf2.paragraphs[0], subtitle, FONT_SUBTITLE, RGBColor(0xCA, 0xDC, 0xFC))

    # 페이지 번호
    txbox3 = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.2), Inches(0.25), Inches(0.8), Inches(0.25)
    )
    tf3 = txbox3.text_frame
    _set_run(tf3.paragraphs[0], page_text, FONT_PAGE, COLORS["text_white"], align=PP_ALIGN.RIGHT)


def add_section_title(slide, left, top, width, text):
    """섹션 제목 (라벨 배경 + 텍스트)"""
    # 배경
    shape = slide.shapes.add_shape(
        1, left, top, width, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()

    # 텍스트
    txbox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.02), width - Inches(0.3), Inches(0.26)
    )
    tf = txbox.text_frame
    _set_run(tf.paragraphs[0], text, FONT_SECTION, COLORS["text_white"], bold=True)


def add_card(slide, left, top, width, height, label, value):
    """라벨-값 카드 (배경 + TextBox 2개)"""
    card_w = width
    label_w = Inches(1.6)
    value_w = card_w - label_w

    # 라벨 배경
    lbl_shape = slide.shapes.add_shape(1, left, top, label_w, height)
    lbl_shape.fill.solid()
    lbl_shape.fill.fore_color.rgb = COLORS["card_label"]
    lbl_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    lbl_shape.line.width = Pt(0.5)

    # 라벨 텍스트
    txbox = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(0.03), label_w - Inches(0.16), height - Inches(0.06)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    _set_run(tf.paragraphs[0], label, FONT_LABEL, COLORS["text_dark"], bold=True)

    # 값 배경
    val_shape = slide.shapes.add_shape(1, left + label_w, top, value_w, height)
    val_shape.fill.solid()
    val_shape.fill.fore_color.rgb = COLORS["card_bg"]
    val_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    val_shape.line.width = Pt(0.5)

    # 값 텍스트
    txbox2 = slide.shapes.add_textbox(
        left + label_w + Inches(0.08), top + Inches(0.03),
        value_w - Inches(0.16), height - Inches(0.06)
    )
    tf2 = txbox2.text_frame
    tf2.word_wrap = True
    _set_run(tf2.paragraphs[0], str(value), FONT_VALUE, COLORS["text_dark"])


def add_bullet_box(slide, left, top, width, height, items):
    """글머리 기호 텍스트 박스"""
    # 배경
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["card_bg"]
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)

    # 텍스트
    txbox = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.05),
        width - Inches(0.24), height - Inches(0.1)
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        _set_run(p, item, FONT_BODY, COLORS["text_dark"])
        p.space_after = Pt(2)


REQUIRED_KEYS = ["title", "site", "date"]  # supplier는 선택 (미입력 시 "미입력" 표시)


def _validate_input(data):
    """필수 키 검증. 누락 시 ValueError."""
    missing = [k for k in REQUIRED_KEYS if not data.get(k)]
    if missing:
        raise ValueError(f"NCR 필수 입력 누락: {', '.join(missing)}")


def _safe_str(value, default="미입력"):
    """None/빈값 → 기본 텍스트"""
    if value is None or (isinstance(value, str) and not value.strip()):
        return default
    return str(value)


def _safe_list(value):
    """None/비리스트 → 빈 리스트, 각 항목 문자열 변환"""
    if not value or not isinstance(value, list):
        return []
    return [str(item) for item in value]


def _add_photo_or_placeholder(slide, left, top, width, height, photo_path):
    """사진 삽입 (비율 유지) 또는 빈 박스 + 안내 텍스트"""
    if photo_path and Path(photo_path).is_file():
        from PIL import Image
        img = Image.open(photo_path)
        img_w, img_h = img.size
        img.close()

        # 비율 유지하며 영역 내 최대 크기 계산
        scale = min(width / img_w, height / img_h)
        final_w = int(img_w * scale)
        final_h = int(img_h * scale)
        # 중앙 정렬
        x_offset = left + (width - final_w) // 2
        y_offset = top + (height - final_h) // 2

        slide.shapes.add_picture(photo_path, x_offset, y_offset, final_w, final_h)
    else:
        # 빈 박스 + "사진 미첨부" 텍스트
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        box.line.width = Pt(1)

        txbox = slide.shapes.add_textbox(
            left + width // 4, top + height // 2 - Inches(0.15),
            width // 2, Inches(0.3)
        )
        tf = txbox.text_frame
        _set_run(tf.paragraphs[0], "사진 미첨부", Pt(14),
                 RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)


def generate_ncr(data, output_path="ncr_output.pptx"):
    """NCR 보고서 PPTX 생성"""
    _validate_input(data)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank

    # ── Slide 1: 기본 정보 ──
    slide1 = prs.slides.add_slide(blank_layout)

    # 배경
    bg = slide1.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg_light"]
    bg.line.fill.background()

    add_header(slide1, "NCR | Non-Conformity Report", _safe_str(data.get("header_subtitle"), "부적합 보고서"), "1 / 2")

    y = HEADER_H + Inches(0.12)
    col1_x = MARGIN
    col1_w = Inches(5.9)
    card_h = Inches(0.32)

    # 섹션: 기본 정보
    add_section_title(slide1, col1_x, y, Inches(12.3), "기본 정보")
    y += Inches(0.38)

    # 좌측 카드들
    cards_left = [
        ("Title", _safe_str(data.get("title"))),
        ("Supplier & Team", _safe_str(data.get("supplier"))),
        ("Occurrence Site", _safe_str(data.get("site"))),
        ("Date", _safe_str(data.get("date"))),
        ("Part No", _safe_str(data.get("part_no"))),
        ("Part Name", _safe_str(data.get("part_name"))),
    ]

    for label, value in cards_left:
        add_card(slide1, col1_x, y, col1_w, card_h, label, value)
        y += card_h + CARD_GAP

    # 우측: 현상 사진 영역 (원본 좌표 기준: L=5.85" T=1.0")
    photo_x = MARGIN + Inches(6.2)
    photo_section_y = HEADER_H + Inches(0.12)
    photo_section_w = Inches(6.1)
    add_section_title(slide1, photo_x, photo_section_y, photo_section_w, "현상 사진")
    photo_y = photo_section_y + Inches(0.38)
    photo_h = Inches(3.2)
    _add_photo_or_placeholder(
        slide1, photo_x, photo_y, photo_section_w, photo_h,
        data.get("photo_path")
    )

    # 우측 하단: 부적합 수량 + NG 표시
    ng_y = photo_y + photo_h + Inches(0.15)
    add_card(slide1, photo_x, ng_y, photo_section_w, card_h, "부적합 수량", _safe_str(data.get("ng_qty")))
    ng_y += card_h + CARD_GAP

    # NG 강조 표시
    ng_box = slide1.shapes.add_shape(1, photo_x, ng_y, Inches(1.2), Inches(0.4))
    ng_box.fill.solid()
    ng_box.fill.fore_color.rgb = COLORS["fail_red"]
    ng_box.line.fill.background()
    ng_txbox = slide1.shapes.add_textbox(photo_x + Inches(0.08), ng_y + Inches(0.03), Inches(1.04), Inches(0.34))
    ng_tf = ng_txbox.text_frame
    _set_run(ng_tf.paragraphs[0], "NG", Pt(24), COLORS["text_white"], bold=True, align=PP_ALIGN.CENTER)

    # 부적합 내용
    desc_y = y + Inches(0.08)
    add_section_title(slide1, col1_x, desc_y, Inches(12.3), "부적합 내용")
    desc_y += Inches(0.35)
    desc_items = _safe_list(data.get("description"))
    if not desc_items:
        desc_items = ["(내용 없음)"]
    add_bullet_box(slide1, col1_x, desc_y, Inches(12.3), Inches(0.7), desc_items)

    # 특기사항
    note_y = desc_y + Inches(0.78)
    add_section_title(slide1, col1_x, note_y, Inches(6.0), "특기사항")
    note_y += Inches(0.35)
    add_card(slide1, col1_x, note_y, Inches(6.0), card_h, "특기사항", _safe_str(data.get("note")))
    note_y += card_h + CARD_GAP
    add_card(slide1, col1_x, note_y, Inches(6.0), card_h, "대책요구일", _safe_str(data.get("action_due")))

    # ── Slide 2: 원인 · 대책 ──
    slide2 = prs.slides.add_slide(blank_layout)

    bg2 = slide2.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = COLORS["bg_light"]
    bg2.line.fill.background()

    add_header(slide2, "NCR | Cause & Action Plan", "원인 · 즉시조치 · 재발방지대책", "2 / 2")

    y2 = HEADER_H + Inches(0.1)
    full_w = Inches(12.3)

    # 원인 분석
    add_section_title(slide2, MARGIN, y2, full_w, "원인 분석")
    y2 += Inches(0.35)

    # 직접 원인
    txbox_lbl = slide2.shapes.add_textbox(MARGIN + Inches(0.1), y2, Inches(2), Inches(0.22))
    tf_lbl = txbox_lbl.text_frame
    _set_run(tf_lbl.paragraphs[0], "① 직접 원인", FONT_LABEL, COLORS["primary"], bold=True)
    y2 += Inches(0.24)

    dc = _safe_list(data.get("direct_cause")) or ["(원인 미입력)"]
    add_bullet_box(slide2, MARGIN, y2, full_w, Inches(0.5), dc)
    y2 += Inches(0.55)

    # 관리 원인
    txbox_lbl2 = slide2.shapes.add_textbox(MARGIN + Inches(0.1), y2, Inches(2), Inches(0.22))
    tf_lbl2 = txbox_lbl2.text_frame
    _set_run(tf_lbl2.paragraphs[0], "② 관리 원인", FONT_LABEL, COLORS["primary"], bold=True)
    y2 += Inches(0.24)

    mc = _safe_list(data.get("mgmt_cause")) or ["(원인 미입력)"]
    add_bullet_box(slide2, MARGIN, y2, full_w, Inches(0.5), mc)
    y2 += Inches(0.58)

    # 즉시 조치
    add_section_title(slide2, MARGIN, y2, full_w, "즉시 조치")
    y2 += Inches(0.35)
    ia = _safe_list(data.get("immediate_action")) or ["(조치 미입력)"]
    add_bullet_box(slide2, MARGIN, y2, full_w, Inches(0.65), ia)
    y2 += Inches(0.72)

    # 재발 방지 대책
    add_section_title(slide2, MARGIN, y2, full_w, "재발 방지 대책")
    y2 += Inches(0.35)
    pv = _safe_list(data.get("prevention")) or ["(대책 미입력)"]
    add_bullet_box(slide2, MARGIN, y2, full_w, Inches(0.65), pv)
    y2 += Inches(0.72)

    # 실행 관리
    add_section_title(slide2, MARGIN, y2, full_w, "실행 관리")
    y2 += Inches(0.35)
    exec_cards = [
        ("담당", _safe_str(data.get("responsible"))),
        ("완료예정", _safe_str(data.get("due_date"))),
        ("비고", _safe_str(data.get("remark"))),
    ]
    for label, value in exec_cards:
        add_card(slide2, MARGIN, y2, Inches(4.0), card_h, label, value)
        y2 += card_h + CARD_GAP

    prs.save(output_path)
    return output_path


# ── 샘플 데이터 (각인텅 대책서 재현) ──
SAMPLE_DATA = {
    "header_subtitle": "각인텅 무단 폐기 · 보고",
    "photo_path": None,  # 사진 경로 (없으면 "사진 미첨부" 표시)
    "title": "각인텅 (Engraving Tong)",
    "supplier": "대원테크 / 팀",
    "site": "SD9A01 라인",
    "date": "2026.03.30",
    "part_no": "각인텅 품번 [확인 필요]",
    "part_name": "각인텅 (Engraving Tong)",
    "ng_qty": "200 EA",
    "description": [
        "생산 공정 후 잔여 각인텅을 불용품(GERP) 등록 없이 임의 폐기",
        "GERP 불용품 등록 · 승인 절차 미준수",
    ],
    "note": "각인텅 무단 폐기 [절차 미준수]",
    "action_due": "2026.3.30",
    "direct_cause": [
        "생산 공정 후 잔여 각인텅을 불용품 등록 없이 임의 폐기",
        "불용품 등록 및 승인 절차 미준수",
    ],
    "mgmt_cause": [
        "불용품 등록 절차 번거로움에 대한 현장 인식",
        "재고 실사 주기 및 관리자 확인 체계 미흡",
    ],
    "immediate_action": [
        "1) 각인텅 재고 전수 실사 및 현황 재확인",
        "2) 관련 담당자 경위 확인 및 내부 조치 기준 검토",
        "3) 전 관리자 대상 불용품 등록 / 승인 절차 재교육",
    ],
    "prevention": [
        "1) 불용품 발생 시 ERP 등록 → 관리자 승인 → 폐기 절차 의무화",
        "2) 월 1회 재고 실사 및 불용품 관리 점검 시행",
        "3) 임의 폐기 적발 시 책임 기준 명확화 및 재교육 연계",
    ],
    "responsible": "현장관리자",
    "due_date": "즉시 적용",
    "remark": "",
}


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="NCR PPT Generator")
    parser.add_argument("--input", "-i", help="JSON 입력 파일 경로")
    parser.add_argument("--output", "-o", default="ncr_output.pptx", help="출력 PPTX 경로")
    args = parser.parse_args()

    if args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = SAMPLE_DATA
        print("샘플 데이터로 생성합니다...")

    out = generate_ncr(data, args.output)
    print(f"생성 완료: {out}")
