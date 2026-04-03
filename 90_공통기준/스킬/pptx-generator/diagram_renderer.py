"""
다이어그램 렌더러 — Graphviz 기반 순서도/프로세스도 생성 + PPTX 삽입
GPT 합의: SKILL.md Visualize 엔진 C층 확장

사용법:
  from diagram_renderer import render_flowchart, render_process, insert_diagram

  # 순서도 생성
  png_path = render_flowchart(
      nodes=[("A", "원재료 입고"), ("B", "수입검사"), ("C", "합격?"), ("D", "입고"), ("E", "반품")],
      edges=[("A","B"), ("B","C"), ("C","D","합격"), ("C","E","불합격")],
      title="수입검사 절차"
  )

  # PPTX 슬라이드에 삽입
  insert_diagram(slide, png_path, left, top, max_width, max_height)

의존성: graphviz (Python), Graphviz 바이너리 (dot)
"""

import os
import sys
import tempfile
from pathlib import Path

try:
    import graphviz
    HAS_GRAPHVIZ = True
except ImportError:
    HAS_GRAPHVIZ = False

from pptx.util import Inches, Emu


# ── Beautify 색상 (SKILL.md 팔레트 기준) ──
DIAGRAM_COLORS = {
    "node_fill":     "#1E2761",  # 진한 남색
    "node_font":     "#FFFFFF",
    "decision_fill": "#E8702A",  # 주황
    "decision_font": "#FFFFFF",
    "start_fill":    "#2C5F2D",  # 녹색
    "start_font":    "#FFFFFF",
    "end_fill":      "#C0392B",  # 적색
    "end_font":      "#FFFFFF",
    "edge_color":    "#333333",
    "edge_font":     "#555555",
    "bg_color":      "#F2F2F2",
}

# ── 노드 유형 자동 판정 키워드 ──
START_KEYWORDS = ["시작", "start", "입고", "접수", "begin"]
END_KEYWORDS = ["종료", "end", "완료", "폐기", "반품", "finish"]
DECISION_KEYWORDS = ["?", "판정", "검사", "확인", "합격", "여부", "분기"]


def _detect_node_type(label: str) -> str:
    """노드 라벨에서 유형 자동 판정"""
    label_lower = label.lower().strip()
    if any(kw in label_lower for kw in DECISION_KEYWORDS):
        return "decision"
    if any(kw in label_lower for kw in START_KEYWORDS):
        return "start"
    if any(kw in label_lower for kw in END_KEYWORDS):
        return "end"
    return "process"


def _apply_node_style(graph, node_id: str, label: str, node_type: str = None):
    """노드 유형에 따른 스타일 적용"""
    if node_type is None:
        node_type = _detect_node_type(label)

    if node_type == "decision":
        graph.node(node_id, label,
                   shape="diamond",
                   style="filled",
                   fillcolor=DIAGRAM_COLORS["decision_fill"],
                   fontcolor=DIAGRAM_COLORS["decision_font"],
                   fontsize="12",
                   width="1.8",
                   height="1.0")
    elif node_type == "start":
        graph.node(node_id, label,
                   shape="oval",
                   style="filled",
                   fillcolor=DIAGRAM_COLORS["start_fill"],
                   fontcolor=DIAGRAM_COLORS["start_font"],
                   fontsize="12")
    elif node_type == "end":
        graph.node(node_id, label,
                   shape="oval",
                   style="filled",
                   fillcolor=DIAGRAM_COLORS["end_fill"],
                   fontcolor=DIAGRAM_COLORS["end_font"],
                   fontsize="12")
    else:  # process
        graph.node(node_id, label,
                   shape="box",
                   style="filled,rounded",
                   fillcolor=DIAGRAM_COLORS["node_fill"],
                   fontcolor=DIAGRAM_COLORS["node_font"],
                   fontsize="12",
                   margin="0.15,0.08")


def render_flowchart(nodes, edges, title=None, direction="TB", output_path=None, dpi=200):
    """
    순서도(Flowchart) 렌더링

    Args:
        nodes: [(id, label), ...] 또는 [(id, label, type), ...]
               type: "start" | "end" | "decision" | "process" (생략 시 자동 판정)
        edges: [(from, to), ...] 또는 [(from, to, label), ...]
        title: 다이어그램 제목 (상단 표시)
        direction: "TB" (위→아래), "LR" (좌→우)
        output_path: 출력 PNG 경로 (None이면 임시 파일)
        dpi: 해상도 (기본 200)

    Returns:
        str: 생성된 PNG 파일 경로
    """
    if not HAS_GRAPHVIZ:
        raise RuntimeError("graphviz 패키지가 설치되지 않았습니다: pip install graphviz")

    # Graphviz 바이너리 PATH 보장
    _ensure_graphviz_path()

    g = graphviz.Digraph(format="png")
    g.attr(rankdir=direction, bgcolor=DIAGRAM_COLORS["bg_color"],
           dpi=str(dpi), pad="0.3", nodesep="0.4", ranksep="0.5")
    g.attr("node", fontname="Malgun Gothic")
    g.attr("edge", fontname="Malgun Gothic", fontsize="10",
           color=DIAGRAM_COLORS["edge_color"],
           fontcolor=DIAGRAM_COLORS["edge_font"])

    # 제목
    if title:
        g.attr(label=title, labelloc="t", fontsize="16",
               fontname="Malgun Gothic", fontcolor=DIAGRAM_COLORS["node_fill"])

    # 노드
    for node in nodes:
        if len(node) == 3:
            nid, label, ntype = node
        else:
            nid, label = node
            ntype = None
        _apply_node_style(g, nid, label, ntype)

    # 엣지
    for edge in edges:
        if len(edge) == 3:
            src, dst, elabel = edge
            g.edge(src, dst, label=f" {elabel} ")
        else:
            src, dst = edge
            g.edge(src, dst)

    # 출력 경로
    if output_path is None:
        tmp = tempfile.NamedTemporaryFile(suffix="", prefix="diagram_", delete=False,
                                          dir=tempfile.gettempdir())
        tmp.close()
        output_path = tmp.name

    result = g.render(output_path, cleanup=True)
    return result


def render_process(steps, title=None, direction="LR", output_path=None, dpi=200):
    """
    프로세스 흐름도 (단순 순차) 렌더링

    Args:
        steps: ["단계1", "단계2", ...] 문자열 리스트
        title: 다이어그램 제목
        direction: "LR" (기본 좌→우) 또는 "TB"
        output_path: 출력 PNG 경로
        dpi: 해상도

    Returns:
        str: 생성된 PNG 파일 경로
    """
    nodes = []
    edges = []
    for i, step in enumerate(steps):
        nid = f"S{i}"
        if i == 0:
            ntype = "start"
        elif i == len(steps) - 1:
            ntype = "end"
        else:
            ntype = "process"
        nodes.append((nid, step, ntype))
        if i > 0:
            edges.append((f"S{i-1}", nid))

    return render_flowchart(nodes, edges, title=title, direction=direction,
                            output_path=output_path, dpi=dpi)


def render_org_chart(data, title=None, output_path=None, dpi=200):
    """
    조직도/계층도 렌더링

    Args:
        data: {"root": "공장장", "children": [
                 {"name": "생산팀장", "children": [{"name": "A라인"}, {"name": "B라인"}]},
                 {"name": "품질팀장", "children": [{"name": "검사반"}]}
               ]}
        title: 다이어그램 제목
        output_path: 출력 PNG 경로
        dpi: 해상도

    Returns:
        str: 생성된 PNG 파일 경로
    """
    if not HAS_GRAPHVIZ:
        raise RuntimeError("graphviz 패키지가 설치되지 않았습니다")

    _ensure_graphviz_path()

    g = graphviz.Digraph(format="png")
    g.attr(rankdir="TB", bgcolor=DIAGRAM_COLORS["bg_color"],
           dpi=str(dpi), pad="0.3", nodesep="0.6", ranksep="0.5")
    g.attr("node", fontname="Malgun Gothic", shape="box", style="filled,rounded",
           fillcolor=DIAGRAM_COLORS["node_fill"], fontcolor=DIAGRAM_COLORS["node_font"],
           fontsize="12", margin="0.15,0.08")
    g.attr("edge", color=DIAGRAM_COLORS["edge_color"], arrowsize="0.7")

    if title:
        g.attr(label=title, labelloc="t", fontsize="16",
               fontname="Malgun Gothic", fontcolor=DIAGRAM_COLORS["node_fill"])

    counter = [0]

    def _add_tree(parent_id, node_data):
        children = node_data.get("children", [])
        for child in children:
            counter[0] += 1
            cid = f"N{counter[0]}"
            g.node(cid, child["name"])
            g.edge(parent_id, cid)
            _add_tree(cid, child)

    root_name = data.get("root", "Root")
    g.node("N0", root_name, fillcolor=DIAGRAM_COLORS["decision_fill"])
    _add_tree("N0", data)

    if output_path is None:
        tmp = tempfile.NamedTemporaryFile(suffix="", prefix="orgchart_", delete=False,
                                          dir=tempfile.gettempdir())
        tmp.close()
        output_path = tmp.name

    return g.render(output_path, cleanup=True)


def insert_diagram(slide, png_path, left, top, max_width, max_height):
    """
    PNG 다이어그램을 PPTX 슬라이드에 비율 유지 삽입

    Args:
        slide: python-pptx Slide 객체
        png_path: PNG 파일 경로
        left, top: 삽입 위치 (Emu/Inches)
        max_width, max_height: 최대 크기 (Emu/Inches)
    """
    from PIL import Image

    img = Image.open(png_path)
    img_w, img_h = img.size
    img.close()

    # Emu로 변환
    if isinstance(max_width, (int, float)) and max_width > 100000:
        mw = max_width
    else:
        mw = int(max_width)
    if isinstance(max_height, (int, float)) and max_height > 100000:
        mh = max_height
    else:
        mh = int(max_height)

    # 비율 유지 스케일
    scale = min(mw / img_w, mh / img_h)
    final_w = int(img_w * scale)
    final_h = int(img_h * scale)

    # 영역 내 중앙 정렬
    x = int(left) + (mw - final_w) // 2
    y = int(top) + (mh - final_h) // 2

    slide.shapes.add_picture(png_path, x, y, final_w, final_h)


def _ensure_graphviz_path():
    """Windows Graphviz 설치 경로를 PATH에 추가"""
    graphviz_paths = [
        r"C:\Program Files\Graphviz\bin",
        r"C:\Program Files (x86)\Graphviz\bin",
    ]
    for p in graphviz_paths:
        if os.path.isdir(p) and p not in os.environ.get("PATH", ""):
            os.environ["PATH"] = os.environ.get("PATH", "") + os.pathsep + p
            break


# ── 시각 타입 자동 선택기 (SKILL.md B층) ──
def auto_select_visual(data_description: str) -> str:
    """
    데이터 특성 설명에서 최적 시각 타입을 추천

    Args:
        data_description: "원인-대책 쌍", "프로세스 절차", "수치 비교" 등

    Returns:
        str: "flowchart" | "process" | "org_chart" | "bar_chart" | "pie_chart" | "table" | "kpi_card"
    """
    desc = data_description.lower()

    flowchart_kw = ["순서도", "흐름도", "flowchart", "워크플로우", "workflow", "분기", "판정"]
    process_kw = ["프로세스", "절차", "단계", "process", "step", "순서"]
    org_kw = ["조직도", "계층", "트리", "org", "hierarchy", "체계"]
    bar_kw = ["비교", "막대", "bar", "추이", "월별", "라인별"]
    pie_kw = ["비율", "구성", "pie", "점유", "비중"]
    table_kw = ["원인", "대책", "표", "table", "목록", "리스트"]
    kpi_kw = ["kpi", "핵심", "달성률", "총합", "요약"]

    for kw_list, vtype in [
        (flowchart_kw, "flowchart"),
        (process_kw, "process"),
        (org_kw, "org_chart"),
        (bar_kw, "bar_chart"),
        (pie_kw, "pie_chart"),
        (kpi_kw, "kpi_card"),
        (table_kw, "table"),
    ]:
        if any(kw in desc for kw in kw_list):
            return vtype

    return "table"  # 기본값


if __name__ == "__main__":
    """테스트 실행"""
    _ensure_graphviz_path()

    print("=== diagram_renderer 테스트 ===\n")

    # 1. 순서도 테스트
    print("1. 순서도 (수입검사 절차)")
    path1 = render_flowchart(
        nodes=[
            ("A", "원재료 입고"),
            ("B", "수입검사"),
            ("C", "합격?"),
            ("D", "입고 처리"),
            ("E", "반품 처리"),
        ],
        edges=[
            ("A", "B"),
            ("B", "C"),
            ("C", "D", "합격"),
            ("C", "E", "불합격"),
        ],
        title="수입검사 절차"
    )
    print(f"   생성: {path1} ({os.path.getsize(path1):,} bytes)")

    # 2. 프로세스 흐름도 테스트
    print("2. 프로세스 흐름도 (불용품 절차)")
    path2 = render_process(
        steps=["불용품 발생", "ERP 등록", "관리자 승인", "폐기 처리", "완료"],
        title="불용품 관리 절차"
    )
    print(f"   생성: {path2} ({os.path.getsize(path2):,} bytes)")

    # 3. 조직도 테스트
    print("3. 조직도 (생산조직)")
    path3 = render_org_chart(
        data={
            "root": "공장장",
            "children": [
                {"name": "생산팀장", "children": [
                    {"name": "SD9A01"},
                    {"name": "SD9A02"},
                ]},
                {"name": "품질팀장", "children": [
                    {"name": "검사반"},
                    {"name": "시정조치"},
                ]},
            ]
        },
        title="생산 조직도"
    )
    print(f"   생성: {path3} ({os.path.getsize(path3):,} bytes)")

    # 4. 시각 타입 자동 선택 테스트
    print("\n4. 시각 타입 자동 선택")
    tests = [
        "수입검사 순서도", "불용품 처리 절차", "생산 조직도",
        "라인별 생산량 비교", "불량 비율 구성", "원인-대책 쌍", "월간 KPI 요약"
    ]
    for t in tests:
        print(f"   '{t}' → {auto_select_visual(t)}")

    # 5. PPTX 삽입 통합 테스트
    print("\n5. PPTX 삽입 통합 테스트")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 순서도 슬라이드
    s1 = prs.slides.add_slide(blank)
    txbox = s1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
    tf = txbox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "수입검사 절차 순서도"
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
    run.font.bold = True
    insert_diagram(s1, path1, Inches(1), Inches(1), Inches(11), Inches(6))

    # 프로세스 슬라이드
    s2 = prs.slides.add_slide(blank)
    txbox2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
    tf2 = txbox2.text_frame
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "불용품 관리 프로세스"
    run2.font.size = Pt(28)
    run2.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
    run2.font.bold = True
    insert_diagram(s2, path2, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))

    # 조직도 슬라이드
    s3 = prs.slides.add_slide(blank)
    txbox3 = s3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
    tf3 = txbox3.text_frame
    run3 = tf3.paragraphs[0].add_run()
    run3.text = "생산 조직도"
    run3.font.size = Pt(28)
    run3.font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
    run3.font.bold = True
    insert_diagram(s3, path3, Inches(2), Inches(1), Inches(9), Inches(6))

    out_path = os.path.join(os.path.dirname(__file__), "diagram_test_output.pptx")
    prs.save(out_path)
    print(f"   PPTX 저장: {out_path}")

    # cleanup
    for p in [path1, path2, path3]:
        try:
            os.remove(p)
        except:
            pass

    print("\n=== ALL TESTS PASSED ===")
