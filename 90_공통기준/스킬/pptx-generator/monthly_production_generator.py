"""
월간 생산실적 PPT 생성기 — PoC v1
GPT 합의: 요약 1장 + 실적표 1장 + 차트 1~2장

입력: 생산실적 데이터 dict (라인별 생산량/목표/전월)
출력: 3~4슬라이드 PPTX

사용법:
  PYTHONUTF8=1 python monthly_production_generator.py
  PYTHONUTF8=1 python monthly_production_generator.py --input data.json --output report.pptx
"""

import json
import sys
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm


# ── 한글 폰트 설정 ──
def setup_korean_font():
    """시스템에서 한글 폰트를 찾아 matplotlib에 설정"""
    korean_fonts = ['Malgun Gothic', 'NanumGothic', 'AppleGothic', 'Gulim']
    for font_name in korean_fonts:
        matches = fm.findSystemFonts(fontpaths=None)
        for fp in matches:
            try:
                f = fm.FontProperties(fname=fp)
                if font_name.lower() in f.get_name().lower():
                    plt.rcParams['font.family'] = f.get_name()
                    plt.rcParams['axes.unicode_minus'] = False
                    return f.get_name()
            except Exception:
                continue
    return None

setup_korean_font()


# ── Beautify 색상 팔레트 ──
COLORS = {
    "primary":    RGBColor(0x1E, 0x27, 0x61),
    "bg_light":   RGBColor(0xF2, 0xF2, 0xF2),
    "accent":     RGBColor(0xE8, 0x70, 0x2A),
    "text_dark":  RGBColor(0x33, 0x33, 0x33),
    "text_white": RGBColor(0xFF, 0xFF, 0xFF),
    "card_bg":    RGBColor(0xFF, 0xFF, 0xFF),
    "pass_green": RGBColor(0x2C, 0x5F, 0x2D),
    "fail_red":   RGBColor(0xC0, 0x39, 0x2B),
}

# matplotlib 색상
CHART_COLORS = {
    "target": "#AAAAAA",
    "prev":   "#7BAFD4",
    "actual": "#1E2761",
    "accent": "#E8702A",
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
HEADER_H = Inches(0.8)


def _set_run(paragraph, text, font_size, color, bold=False, align=None):
    """run 단위 텍스트 설정"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    if align is not None:
        paragraph.alignment = align


def add_header(slide, title, subtitle, page_text):
    """헤더 바"""
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, HEADER_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary"]
    shape.line.fill.background()

    txbox = slide.shapes.add_textbox(MARGIN, Inches(0.08), Inches(8), Inches(0.4))
    tf = txbox.text_frame
    tf.word_wrap = True
    _set_run(tf.paragraphs[0], title, Pt(28), COLORS["text_white"], bold=True)

    txbox2 = slide.shapes.add_textbox(MARGIN, Inches(0.48), Inches(10), Inches(0.25))
    tf2 = txbox2.text_frame
    _set_run(tf2.paragraphs[0], subtitle, Pt(14), RGBColor(0xCA, 0xDC, 0xFC))

    txbox3 = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), Inches(0.25), Inches(0.8), Inches(0.25))
    tf3 = txbox3.text_frame
    _set_run(tf3.paragraphs[0], page_text, Pt(10), COLORS["text_white"], align=PP_ALIGN.RIGHT)


def add_kpi_card(slide, left, top, width, height, label, value, sub_text="", color=None):
    """KPI 카드 (큰 숫자 + 라벨)"""
    # 배경
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["card_bg"]
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)

    # 상단 악센트 바
    accent = slide.shapes.add_shape(1, left, top, width, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color or COLORS["primary"]
    accent.line.fill.background()

    # 값 (큰 숫자)
    txbox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), Inches(0.6))
    tf = txbox.text_frame
    _set_run(tf.paragraphs[0], str(value), Pt(36), color or COLORS["primary"], bold=True, align=PP_ALIGN.CENTER)

    # 라벨
    txbox2 = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.8), width - Inches(0.2), Inches(0.3))
    tf2 = txbox2.text_frame
    _set_run(tf2.paragraphs[0], label, Pt(11), COLORS["text_dark"], align=PP_ALIGN.CENTER)

    # 서브텍스트
    if sub_text:
        txbox3 = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.05), width - Inches(0.2), Inches(0.2))
        tf3 = txbox3.text_frame
        sub_color = COLORS["pass_green"] if "+" not in sub_text and "▲" not in sub_text else COLORS["fail_red"]
        if "▲" in sub_text or "+" in sub_text:
            sub_color = COLORS["pass_green"]
        if "▼" in sub_text or "-" in sub_text:
            sub_color = COLORS["fail_red"]
        _set_run(tf3.paragraphs[0], sub_text, Pt(10), sub_color, align=PP_ALIGN.CENTER)


def create_bar_chart(data, output_path):
    """막대 차트 생성 (목표/전월/실적)"""
    lines = data.get("lines", [])
    labels = [d["name"] for d in lines]
    targets = [d.get("target", 0) for d in lines]
    prev = [d.get("prev_month", 0) for d in lines]
    actual = [d.get("actual", 0) for d in lines]

    fig, ax = plt.subplots(figsize=(11, 4.5))

    x = range(len(labels))
    bar_w = 0.25
    bars1 = ax.bar([i - bar_w for i in x], targets, bar_w, label='목표',
                   color=CHART_COLORS["target"], edgecolor='none')
    bars2 = ax.bar(x, prev, bar_w, label='전월',
                   color=CHART_COLORS["prev"], edgecolor='none')
    bars3 = ax.bar([i + bar_w for i in x], actual, bar_w, label='실적',
                   color=CHART_COLORS["actual"], edgecolor='none')

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=10)
    ax.set_ylabel('수량', fontsize=11)
    ax.legend(fontsize=10, loc='upper right')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)

    # 실적 데이터 라벨
    for bar in bars3:
        height = bar.get_height()
        if height > 0:
            ax.text(bar.get_x() + bar.get_width()/2., height + max(actual)*0.01,
                    f'{int(height):,}', ha='center', va='bottom', fontsize=8,
                    color=CHART_COLORS["actual"])

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return output_path


def create_achievement_chart(data, output_path):
    """달성률 가로 막대 차트"""
    lines = data.get("lines", [])
    labels = [d["name"] for d in lines]
    rates = []
    for d in lines:
        target = d.get("target", 1)
        actual = d.get("actual", 0)
        rates.append(round(actual / target * 100, 1) if target > 0 else 0)

    fig, ax = plt.subplots(figsize=(11, 4))

    colors = [CHART_COLORS["actual"] if r >= 100 else CHART_COLORS["accent"] for r in rates]
    bars = ax.barh(labels, rates, color=colors, edgecolor='none', height=0.6)

    ax.axvline(x=100, color='#C0392B', linestyle='--', linewidth=1, alpha=0.7)
    ax.set_xlabel('달성률 (%)', fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='x', alpha=0.3)

    for bar, rate in zip(bars, rates):
        ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height()/2,
                f'{rate}%', ha='left', va='center', fontsize=10, fontweight='bold')

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return output_path


REQUIRED_KEYS = ["month", "lines"]


def _validate_input(data):
    """필수 키 검증."""
    missing = [k for k in REQUIRED_KEYS if not data.get(k)]
    if missing:
        raise ValueError(f"월간 생산실적 필수 입력 누락: {', '.join(missing)}")
    lines = data["lines"]
    if not isinstance(lines, list) or len(lines) == 0:
        raise ValueError("lines는 1개 이상의 라인 데이터가 필요합니다")
    for i, line in enumerate(lines):
        if not line.get("name"):
            raise ValueError(f"lines[{i}]: name 누락")
        if "actual" not in line:
            raise ValueError(f"lines[{i}] ({line['name']}): actual 누락")


def generate_monthly_report(data, output_path="monthly_production_output.pptx"):
    """월간 생산실적 보고서 PPTX 생성"""
    _validate_input(data)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    month = data.get("month", "2026년 3월")
    lines = data.get("lines", [])

    # 총합 계산
    total_target = sum(d.get("target", 0) for d in lines)
    total_actual = sum(d.get("actual", 0) for d in lines)
    total_prev = sum(d.get("prev_month", 0) for d in lines)
    total_rate = round(total_actual / total_target * 100, 1) if total_target > 0 else 0
    mom_change = total_actual - total_prev
    mom_pct = round(mom_change / total_prev * 100, 1) if total_prev > 0 else 0

    total_slides = 4

    # ── Slide 1: 요약 ──
    slide1 = prs.slides.add_slide(blank_layout)
    bg = slide1.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg_light"]
    bg.line.fill.background()

    add_header(slide1, f"월간 생산실적 — {month}", "Monthly Production Report", f"1 / {total_slides}")

    # KPI 카드 4개
    card_w = Inches(2.8)
    card_h = Inches(1.35)
    card_y = HEADER_H + Inches(0.2)
    gap = Inches(0.3)
    start_x = MARGIN + Inches(0.2)

    add_kpi_card(slide1, start_x, card_y, card_w, card_h,
                 "총 생산량", f"{total_actual:,}", f"목표 대비 {total_rate}%",
                 COLORS["primary"])

    add_kpi_card(slide1, start_x + card_w + gap, card_y, card_w, card_h,
                 "목표", f"{total_target:,}", "",
                 RGBColor(0x88, 0x88, 0x88))

    mom_sign = "▲" if mom_change >= 0 else "▼"
    add_kpi_card(slide1, start_x + (card_w + gap) * 2, card_y, card_w, card_h,
                 "전월 대비", f"{mom_sign} {abs(mom_pct)}%",
                 f"{mom_sign} {abs(mom_change):,}개",
                 COLORS["pass_green"] if mom_change >= 0 else COLORS["fail_red"])

    # 달성률
    rate_color = COLORS["pass_green"] if total_rate >= 100 else COLORS["fail_red"]
    add_kpi_card(slide1, start_x + (card_w + gap) * 3, card_y, card_w, card_h,
                 "달성률", f"{total_rate}%", "목표 100% 기준",
                 rate_color)

    # 요약 텍스트
    summary_y = card_y + card_h + Inches(0.3)
    summary_box = slide1.shapes.add_shape(1, MARGIN, summary_y, Inches(12.3), Inches(3.5))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = COLORS["card_bg"]
    summary_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    summary_box.line.width = Pt(0.5)

    txbox = slide1.shapes.add_textbox(
        MARGIN + Inches(0.2), summary_y + Inches(0.1),
        Inches(11.9), Inches(3.3)
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    # AI 요약 문안 자동 생성
    _set_run(tf.paragraphs[0], f"■ {month} 생산실적 요약", Pt(14), COLORS["primary"], bold=True)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    rate_text = "달성" if total_rate >= 100 else "미달"
    _set_run(p2, f"총 생산량 {total_actual:,}개로 목표 {total_target:,}개 대비 {total_rate}% {rate_text}.",
             Pt(12), COLORS["text_dark"])

    p3 = tf.add_paragraph()
    p3.space_before = Pt(4)
    mom_text = "증가" if mom_change >= 0 else "감소"
    _set_run(p3, f"전월({total_prev:,}개) 대비 {abs(mom_change):,}개 {mom_text} ({mom_sign}{abs(mom_pct)}%).",
             Pt(12), COLORS["text_dark"])

    # 라인별 요약
    p4 = tf.add_paragraph()
    p4.space_before = Pt(12)
    _set_run(p4, "■ 라인별 현황", Pt(14), COLORS["primary"], bold=True)

    for d in lines:
        line_rate = round(d["actual"] / d["target"] * 100, 1) if d.get("target", 0) > 0 else 0
        status = "달성" if line_rate >= 100 else "미달"
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        _set_run(p, f"  {d['name']}: {d['actual']:,}개 / 목표 {d['target']:,}개 ({line_rate}% {status})",
                 Pt(11), COLORS["text_dark"])

    # ── Slide 2: 실적표 ──
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = COLORS["bg_light"]
    bg2.line.fill.background()

    add_header(slide2, f"라인별 상세 실적 — {month}", "Line-by-Line Detail", f"2 / {total_slides}")

    # 테이블 생성
    table_y = HEADER_H + Inches(0.2)
    cols = 6
    rows = len(lines) + 2  # 헤더 + 데이터 + 합계
    table_shape = slide2.shapes.add_table(
        rows, cols, MARGIN, table_y, Inches(12.3), Inches(0.4 * rows)
    )
    table = table_shape.table

    # 열 너비
    col_widths = [Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(1.9), Inches(1.9)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 헤더
    headers = ["라인명", "목표", "실적", "전월", "달성률", "전월 대비"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLORS["text_white"]
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary"]

    # 데이터 행
    for row_idx, d in enumerate(lines):
        row = row_idx + 1
        line_rate = round(d["actual"] / d["target"] * 100, 1) if d.get("target", 0) > 0 else 0
        mom = d["actual"] - d.get("prev_month", 0)
        mom_s = f"{'▲' if mom >= 0 else '▼'} {abs(mom):,}"

        values = [
            d["name"],
            f"{d['target']:,}",
            f"{d['actual']:,}",
            f"{d.get('prev_month', 0):,}",
            f"{line_rate}%",
            mom_s
        ]
        for col_idx, val in enumerate(values):
            cell = table.cell(row, col_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = COLORS["text_dark"]
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if row_idx % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)

    # 합계 행
    total_row = len(lines) + 1
    total_mom = total_actual - total_prev
    totals = [
        "합계",
        f"{total_target:,}",
        f"{total_actual:,}",
        f"{total_prev:,}",
        f"{total_rate}%",
        f"{'▲' if total_mom >= 0 else '▼'} {abs(total_mom):,}"
    ]
    for col_idx, val in enumerate(totals):
        cell = table.cell(total_row, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLORS["primary"]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEB, 0xF0)

    # ── Slide 3: 생산량 비교 차트 ──
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = COLORS["bg_light"]
    bg3.line.fill.background()

    add_header(slide3, f"라인별 생산량 비교 — {month}", "목표 vs 전월 vs 실적", f"3 / {total_slides}")

    chart_path = os.path.join(os.path.dirname(output_path) or ".", "_tmp_bar_chart.png")
    create_bar_chart(data, chart_path)

    chart_y = HEADER_H + Inches(0.2)
    slide3.shapes.add_picture(chart_path, MARGIN, chart_y, Inches(12.3), Inches(6.0))

    # ── Slide 4: 달성률 차트 ──
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = COLORS["bg_light"]
    bg4.line.fill.background()

    add_header(slide4, f"라인별 달성률 — {month}", "Achievement Rate by Line", f"4 / {total_slides}")

    achieve_path = os.path.join(os.path.dirname(output_path) or ".", "_tmp_achieve_chart.png")
    create_achievement_chart(data, achieve_path)

    slide4.shapes.add_picture(achieve_path, MARGIN, HEADER_H + Inches(0.2), Inches(12.3), Inches(5.5))

    prs.save(output_path)

    # 임시 차트 파일 정리
    for tmp in [chart_path, achieve_path]:
        if os.path.exists(tmp):
            os.remove(tmp)

    return output_path


# ── 샘플 데이터 ──
SAMPLE_DATA = {
    "month": "2026년 3월",
    "lines": [
        {"name": "SD9A01", "target": 15000, "actual": 14800, "prev_month": 14200},
        {"name": "SD9A02", "target": 12000, "actual": 12500, "prev_month": 11800},
        {"name": "WABAS01", "target": 8000, "actual": 8200, "prev_month": 7900},
        {"name": "WAMAS01", "target": 6000, "actual": 5800, "prev_month": 6100},
        {"name": "DRAAS11", "target": 4500, "actual": 4600, "prev_month": 4300},
        {"name": "ANAAS04", "target": 3500, "actual": 3400, "prev_month": 3200},
        {"name": "LTAAS01", "target": 2000, "actual": 2100, "prev_month": 1950},
    ]
}


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Monthly Production Report PPT Generator")
    parser.add_argument("--input", "-i", help="JSON 입력 파일 경로")
    parser.add_argument("--output", "-o", default="monthly_production_output.pptx", help="출력 PPTX 경로")
    args = parser.parse_args()

    if args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = SAMPLE_DATA
        print("샘플 데이터로 생성합니다...")

    out = generate_monthly_report(data, args.output)
    print(f"생성 완료: {out}")
