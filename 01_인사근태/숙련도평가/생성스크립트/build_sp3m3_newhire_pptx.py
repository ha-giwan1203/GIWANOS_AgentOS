# -*- coding: utf-8 -*-
"""
SP3M3 신규 입사자 교육자료 — PPTX 생성 + PDF 변환
- pptx-generator 스킬 패턴 (Beautify 규칙 + 회사 표준 색 팔레트)
- 16:9 슬라이드 7장 (표지 + 공정 6장)
- python-pptx 네이티브 텍스트·도형 + PowerPoint COM PDF 변환
"""
from __future__ import annotations
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))
from create_sp3m3_newhire_v5 import (
    load_processes, load_eval_forms, load_workstd_steps, extract_parts,
    PROC_OVERRIDES,
)

OUTDIR = SCRIPT_DIR.parent / "SP3M3_신규입사자 교육자료"
PPTX_OUT = OUTDIR / "SP3M3_신규입사자 교육자료.pptx"
PDF_OUT = OUTDIR / "SP3M3_신규입사자 교육자료_디자인.pdf"

# 회사 표준 팔레트 (pptx-generator MANUAL)
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x0F, 0x14, 0x35)
ORANGE = RGBColor(0xE8, 0x70, 0x2A)
GRAY_BG = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT = RGBColor(0x95, 0x95, 0x95)
GREEN = RGBColor(0x2C, 0x5F, 0x2D)
GREEN_BG = RGBColor(0xE3, 0xF3, 0xE0)
RED = RGBColor(0xC0, 0x39, 0x2B)
RED_BG = RGBColor(0xFC, 0xE4, 0xE0)
YELLOW_ACCENT = RGBColor(0xFA, 0xCC, 0x15)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 슬라이드 크기 (16:9)
W, H = Inches(13.333), Inches(7.5)


# ─────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────
def add_rect(slide, x, y, cx, cy, fill, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    if not shadow:
        shp.shadow.inherit = False
    return shp


def add_round_rect(slide, x, y, cx, cy, fill, line=None, radius=0.05):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, cx, cy, text, *,
             size=14, bold=False, color=GRAY_TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multiline(slide, x, y, cx, cy, lines, *,
                  size=12, color=GRAY_TEXT, bold=False, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.3, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


# ─────────────────────────────────────────────
# 슬라이드 1 — 표지
# ─────────────────────────────────────────────
def build_cover(prs, procs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 배경 단색 네이비
    add_rect(slide, 0, 0, W, H, NAVY_DARK)

    # 우측 절반 약간 밝은 네이비 패널 (시각 흐름 좌→우)
    add_rect(slide, Inches(7.3), 0, Inches(6.04), H, NAVY)

    # 좌측 세로 노랑 막대 (브랜드 액센트)
    add_rect(slide, 0, 0, Inches(0.2), H, YELLOW_ACCENT)

    # ─── 좌측 50% — 타이틀 영역 ───
    add_text(slide, Inches(0.7), Inches(0.6), Inches(6), Inches(0.4),
             "(유)삼송  ·  SEAT BELT ASSEMBLY",
             size=10, color=GRAY_LIGHT)

    chip = add_round_rect(slide, Inches(0.7), Inches(1.1), Inches(2.3), Inches(0.4),
                          NAVY_DARK, radius=0.5)
    chip.line.color.rgb = YELLOW_ACCENT
    chip.line.width = Pt(1)
    add_text(slide, Inches(0.7), Inches(1.1), Inches(2.3), Inches(0.4),
             "▸ SP3M3 LINE", size=10, bold=True, color=YELLOW_ACCENT,
             align=PP_ALIGN.CENTER)

    # 큰 타이틀 — 좌측 정렬, 두 줄
    add_text(slide, Inches(0.7), Inches(2.4), Inches(6.5), Inches(1.2),
             "신규 입사자", size=48, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, Inches(0.7), Inches(3.4), Inches(6.5), Inches(1.2),
             "교육자료", size=48, bold=True, color=YELLOW_ACCENT)

    # 부제 — 좌측 영역 내
    add_multiline(slide, Inches(0.7), Inches(4.8), Inches(6.3), Inches(1.4), [
        "자동차 안전벨트 조립 라인.",
        "1개의 이종·누락이 차량 1대 리콜로 이어집니다.",
        "빠르게보다 정확하게 — 의심되면 즉시 사수 확인.",
    ], size=12, color=RGBColor(0xCB, 0xD5, 0xE1), line_spacing=1.5)

    # 좌측 하단 정보
    add_text(slide, Inches(0.7), Inches(6.7), Inches(2), Inches(0.25),
             "ISSUED", size=8, color=GRAY_LIGHT, bold=True)
    add_text(slide, Inches(0.7), Inches(6.95), Inches(2), Inches(0.3),
             "2026.05", size=13, color=WHITE, bold=True)

    add_text(slide, Inches(2.7), Inches(6.7), Inches(2), Inches(0.25),
             "REVISION", size=8, color=GRAY_LIGHT, bold=True)
    add_text(slide, Inches(2.7), Inches(6.95), Inches(2), Inches(0.3),
             "Rev.1", size=13, color=WHITE, bold=True)

    add_text(slide, Inches(4.7), Inches(6.7), Inches(2.3), Inches(0.25),
             "OWNER", size=8, color=GRAY_LIGHT, bold=True)
    add_text(slide, Inches(4.7), Inches(6.95), Inches(2.3), Inches(0.3),
             "생산관리", size=13, color=WHITE, bold=True)

    # ─── 우측 50% — CONTENTS 6개 카드 (3x2) ───
    add_text(slide, Inches(7.7), Inches(0.6), Inches(5.3), Inches(0.4),
             "CONTENTS", size=10, color=YELLOW_ACCENT, bold=True)
    add_text(slide, Inches(7.7), Inches(1.0), Inches(5.3), Inches(0.5),
             "공정 6개 안내", size=20, bold=True, color=WHITE)

    # 노랑 가로 줄 (구분선)
    add_rect(slide, Inches(7.7), Inches(1.7), Inches(0.5), Inches(0.04),
             YELLOW_ACCENT)

    # 카드 3x2 그리드
    card_area_x = Inches(7.7)
    card_area_y = Inches(2.1)
    card_w = Inches(1.74)
    card_h = Inches(2.15)
    gap_x = Inches(0.12)
    gap_y = Inches(0.18)

    for i, p in enumerate(procs):
        col = i % 3
        row = i // 3
        x = card_area_x + (card_w + gap_x) * col
        y = card_area_y + (card_h + gap_y) * row

        card = add_round_rect(slide, x, y, card_w, card_h,
                              NAVY_DARK, radius=0.08)
        card.line.color.rgb = RGBColor(0x3B, 0x4D, 0x77)
        card.line.width = Pt(0.5)

        # 카드 좌측 노랑 짧은 막대
        add_rect(slide, x, y, Inches(0.06), card_h, YELLOW_ACCENT)

        # 큰 공정 번호
        add_text(slide, x + Inches(0.2), y + Inches(0.2),
                 card_w - Inches(0.3), Inches(0.8),
                 str(p["no"]), size=32, bold=True, color=YELLOW_ACCENT,
                 anchor=MSO_ANCHOR.TOP)

        # 라벨
        add_text(slide, x + Inches(0.2), y + Inches(1.05),
                 card_w - Inches(0.3), Inches(0.25),
                 "PROCESS", size=7, color=GRAY_LIGHT, bold=True,
                 anchor=MSO_ANCHOR.TOP)

        # 공정명 (자동 줄바꿈)
        name = p["name"]
        # 길면 줄바꿈
        if len(name) > 18:
            # & 기준 줄바꿈
            parts_str = name.split(" & ")
            name_display = "\n".join(parts_str[:2])
            if len(parts_str) > 2:
                name_display += f" & {parts_str[2]}"
        else:
            name_display = name
        add_multiline(slide, x + Inches(0.2), y + Inches(1.3),
                      card_w - Inches(0.3), Inches(0.85),
                      [name_display], size=8.5, color=WHITE,
                      anchor=MSO_ANCHOR.TOP, line_spacing=1.25)


# ─────────────────────────────────────────────
# 슬라이드 2~7 — 공정 페이지
# ─────────────────────────────────────────────
SAFETY = [
    ("일상 점검", "에어압 0.4~0.6 MPa"),
    ("초물표 확인", "차종 변경 시 첫 5개 사수"),
    ("장갑·귀마개", "회전부·압입부 손가락 금지"),
    ("이상 시 중지", "관리자 즉시 보고"),
    ("의심 시 멈춤", "NG 흘려보내지 말 것"),
]


def build_proc(prs, proc, eval_data, ws_data, parts, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ov = PROC_OVERRIDES.get(proc["no"], {})
    parts_label = ov.get("parts_label", "사용 부품")
    skip_id = ov.get("skip_identifier_check", False)

    # 배경
    add_rect(slide, 0, 0, W, H, GRAY_BG)

    # ─── 헤더 (네이비 띠) ───
    add_rect(slide, 0, 0, W, Inches(1.5), NAVY_DARK)
    add_rect(slide, 0, 0, Inches(0.3), Inches(1.5), YELLOW_ACCENT)

    add_text(slide, Inches(0.7), Inches(0.18), Inches(6), Inches(0.3),
             f"SP3M3 LINE   ·   PROCESS",
             size=10, color=GRAY_LIGHT)
    add_text(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.6),
             proc["name"], size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.TOP)
    # 공정 번호 칩
    chip = add_round_rect(slide, Inches(11.0), Inches(0.4), Inches(1.6), Inches(0.55),
                          NAVY_DARK, radius=0.5)
    chip.line.color.rgb = YELLOW_ACCENT
    chip.line.width = Pt(1)
    add_text(slide, Inches(11.0), Inches(0.4), Inches(1.6), Inches(0.55),
             f"공정 {proc['no']}", size=12, bold=True, color=YELLOW_ACCENT,
             align=PP_ALIGN.CENTER)

    # 메타 (요구레벨 / 시트참조 / 목표)
    meta_y = Inches(1.1)
    for i, (k, v) in enumerate([
        ("요구레벨", proc["level"]),
        ("작업표준서 시트", ov.get("ref_sheet", "자동")),
        ("학습 목표", "1개월 내 단독 작업"),
    ]):
        x = Inches(0.7 + i * 4.2)
        add_text(slide, x, meta_y, Inches(1.5), Inches(0.3),
                 k, size=9, color=YELLOW_ACCENT, bold=True)
        add_text(slide, x + Inches(1.5), meta_y, Inches(2.7), Inches(0.3),
                 str(v), size=11, color=WHITE)

    # ─── 본문 영역 ───
    body_top = Inches(1.85)

    # ② 사용 부품 (상단 띠)
    add_round_rect(slide, Inches(0.5), body_top, W - Inches(1), Inches(0.85),
                   WHITE, line=RGBColor(0xE5, 0xE7, 0xEB), radius=0.1)
    # 좌측 라벨 박스
    add_round_rect(slide, Inches(0.7), body_top + Inches(0.18),
                   Inches(2.0), Inches(0.5),
                   NAVY, radius=0.3)
    add_text(slide, Inches(0.7), body_top + Inches(0.18),
             Inches(2.0), Inches(0.5),
             parts_label.split("(")[0].strip(), size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # 부품 칩들
    chip_x = Inches(3.0)
    chip_y = body_top + Inches(0.22)
    for part in parts[:8]:
        chip_w = Inches(max(1.0, 0.18 * len(part) + 0.5))
        add_round_rect(slide, chip_x, chip_y, chip_w, Inches(0.42),
                       RGBColor(0xEF, 0xF6, 0xFF),
                       line=RGBColor(0xBF, 0xDB, 0xFE), radius=0.5)
        add_text(slide, chip_x, chip_y, chip_w, Inches(0.42),
                 part, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        chip_x += chip_w + Inches(0.12)

    # ③ 작업 순서 (좌측 절반) — 꼭 지킬 것 띠(6.2") 안 침범하도록 높이 3.3"
    steps_top = body_top + Inches(1.05)
    steps_h = Inches(3.3)
    add_round_rect(slide, Inches(0.5), steps_top, Inches(6.1), steps_h,
                   WHITE, line=RGBColor(0xE5, 0xE7, 0xEB), radius=0.05)
    add_text(slide, Inches(0.8), steps_top + Inches(0.1),
             Inches(5.5), Inches(0.35),
             "③  작업 순서",
             size=13, bold=True, color=NAVY)

    steps = ws_data.get("steps", [])
    n_steps = min(len(steps), 11)
    # 사용 가능 영역 = 박스 높이 - 제목·여백(0.55)
    avail = steps_h - Inches(0.55)
    step_gap = min(Inches(0.34), avail / max(n_steps, 1))
    badge_size = Inches(0.30) if n_steps <= 8 else Inches(0.26)
    text_size = 10 if n_steps <= 8 else 9
    step_y = steps_top + Inches(0.55)
    for i, s in enumerate(steps[:11], 1):
        is_id = (not skip_id) and i == 1 and "식별표" in s
        badge_color = YELLOW_ACCENT if is_id else NAVY
        text_color = NAVY if is_id else WHITE
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.8), step_y, badge_size, badge_size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()
        badge.shadow.inherit = False
        add_text(slide, Inches(0.8), step_y, badge_size, badge_size,
                 str(i), size=text_size - 1, bold=True, color=text_color,
                 align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.2), step_y - Inches(0.02),
                 Inches(5.35), step_gap,
                 s, size=text_size, color=GRAY_TEXT)
        step_y += step_gap

    # ④ OK/NG (우측 절반) — 동일 높이 3.3"
    okng_top = steps_top
    okng_h = steps_h
    okng_x = Inches(6.85)
    okng_w = Inches(5.95)
    add_round_rect(slide, okng_x, okng_top, okng_w, okng_h,
                   WHITE, line=RGBColor(0xE5, 0xE7, 0xEB), radius=0.05)
    add_text(slide, okng_x + Inches(0.3), okng_top + Inches(0.1),
             okng_w - Inches(0.6), Inches(0.35),
             "④  합격(OK)  vs  불합격(NG)", size=13, bold=True, color=NAVY)

    # OK/NG 행 간격도 동적 — 박스 안에 6행이 들어가게
    okng_y = okng_top + Inches(0.55)
    okng_row_h = (okng_h - Inches(0.6)) / 6
    if eval_data:
        items = eval_data["items"][:6]
        crits = eval_data["criteria"][:6]
        cell_h = okng_row_h - Inches(0.05)
        for i, (it, (q, _u, y)) in enumerate(zip(items, crits), 1):
            add_text(slide, okng_x + Inches(0.2), okng_y,
                     Inches(0.3), cell_h,
                     str(i), size=10, bold=True, color=GRAY_LIGHT,
                     align=PP_ALIGN.CENTER)
            add_text(slide, okng_x + Inches(0.55), okng_y,
                     Inches(2.0), cell_h,
                     (it or "")[:32], size=8.5, color=GRAY_TEXT)
            add_round_rect(slide, okng_x + Inches(2.6), okng_y,
                           Inches(1.6), cell_h, GREEN_BG, radius=0.15)
            add_text(slide, okng_x + Inches(2.6), okng_y,
                     Inches(1.6), cell_h,
                     f"✓ {(q or '')[:18]}", size=9, bold=True, color=GREEN,
                     align=PP_ALIGN.CENTER)
            add_round_rect(slide, okng_x + Inches(4.25), okng_y,
                           Inches(1.6), cell_h, RED_BG, radius=0.15)
            add_text(slide, okng_x + Inches(4.25), okng_y,
                     Inches(1.6), cell_h,
                     f"✗ {(y or '')[:18]}", size=9, bold=True, color=RED,
                     align=PP_ALIGN.CENTER)
            okng_y += okng_row_h

    # ⑤ 꼭 지킬 것 (하단 띠)
    safety_y = Inches(6.2)
    safety_h = Inches(0.95)
    add_round_rect(slide, Inches(0.5), safety_y, W - Inches(1), safety_h,
                   RGBColor(0xFF, 0xFB, 0xEB),
                   line=RGBColor(0xFD, 0xE6, 0x8A), radius=0.05)
    add_text(slide, Inches(0.8), safety_y + Inches(0.1),
             Inches(2.5), Inches(0.3),
             "⑤  꼭 지킬 것",
             size=11, bold=True, color=ORANGE)
    safety_x = Inches(3.3)
    each_w = (W - Inches(4.0)) / 5
    for title, desc in SAFETY:
        # 좌측 점
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     safety_x, safety_y + Inches(0.27),
                                     Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE
        dot.line.fill.background()
        dot.shadow.inherit = False
        # 텍스트 2줄
        add_text(slide, safety_x + Inches(0.2), safety_y + Inches(0.1),
                 each_w - Inches(0.25), Inches(0.3),
                 title, size=10, bold=True, color=GRAY_TEXT)
        add_text(slide, safety_x + Inches(0.2), safety_y + Inches(0.4),
                 each_w - Inches(0.25), Inches(0.4),
                 desc, size=8.5, color=GRAY_LIGHT)
        safety_x += each_w

    # 페이지 푸터
    add_text(slide, Inches(0.5), Inches(7.25), Inches(6), Inches(0.2),
             "사수 사인 ____________     라인장 확인 ____________     날짜 26 / __ / __",
             size=9, color=GRAY_LIGHT)
    add_text(slide, W - Inches(2), Inches(7.25), Inches(1.5), Inches(0.2),
             f"{page_no} / {total}", size=9, bold=True, color=NAVY,
             align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────
def main():
    print("[1/4] 데이터 로드...")
    procs = load_processes()
    ef = load_eval_forms()

    print("[2/4] PPTX 생성 (16:9)...")
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    build_cover(prs, procs)
    total = len(procs) + 1
    for i, p in enumerate(procs, start=2):
        ws_data = load_workstd_steps(p["no"])
        parts = extract_parts(p, ef.get(p["no"]), ws_data)
        build_proc(prs, p, ef.get(p["no"]), ws_data, parts, i, total)

    OUTDIR.mkdir(parents=True, exist_ok=True)
    if PPTX_OUT.exists():
        PPTX_OUT.unlink()
    prs.save(str(PPTX_OUT))
    print(f"  → {PPTX_OUT.name}  ({PPTX_OUT.stat().st_size/1024:,.1f} KB)")

    print("[3/4] PowerPoint COM → PDF 변환...")
    import win32com.client as win32
    ppt = win32.gencache.EnsureDispatch("PowerPoint.Application")
    # ppt.Visible = 1 — 일부 환경에서 필수
    try:
        ppt.Visible = 1
    except Exception:
        pass
    if PDF_OUT.exists():
        PDF_OUT.unlink()
    wb = ppt.Presentations.Open(str(PPTX_OUT.resolve()), WithWindow=False)
    # PpSaveAsFileType = 32 (ppSaveAsPDF)
    wb.SaveAs(str(PDF_OUT.resolve()), 32)
    wb.Close()
    ppt.Quit()

    print(f"[4/4] 완료")
    print(f"  PPTX: {PPTX_OUT}")
    print(f"  PDF:  {PDF_OUT}  ({PDF_OUT.stat().st_size/1024:,.1f} KB)")


if __name__ == "__main__":
    main()
